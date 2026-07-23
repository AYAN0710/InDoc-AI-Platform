from pptx import Presentation

def extract_pptx_text(file_path:str)->str:
    presentation=Presentation(file_path)
    extracted_text=[]
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape,"text"):
                if shape.text.strip():
                    extracted_text.append(shape.text)
    return "\n".join(extracted_text)